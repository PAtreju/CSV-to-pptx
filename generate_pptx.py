import copy
from io import BytesIO
import requests
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Pt
from PIL import Image

def _get_name_pic_pairs(slide):
    pairs, shapes = [], list(slide.shapes)
    for i, shp in enumerate(shapes):
        if (
            shp.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
            and shp.text.strip().startswith("Product Name")
        ):
            pic = next(
                (s for s in shapes[i + 1 :] if s.shape_type == MSO_SHAPE_TYPE.PICTURE),
                None,
            )
            if pic:
                pairs.append((shp, pic))
    return pairs


def _duplicate_slide(prs, src_idx):
    src = prs.slides[src_idx]
    new_slide = prs.slides.add_slide(src.slide_layout)
    for ph in list(new_slide.shapes):
        new_slide.shapes._spTree.remove(ph._element)
    for shp in src.shapes:
        new_slide.shapes._spTree.append(copy.deepcopy(shp._element))

    return new_slide


def _replace_picture(slide, pic_shape, img_bytes):
    left, top, box_w, box_h = (
        pic_shape.left,
        pic_shape.top,
        pic_shape.width,
        pic_shape.height,
    )

    slide.shapes._spTree.remove(pic_shape._element)

    with Image.open(BytesIO(img_bytes)) as im:
        iw, ih = im.size
    img_ratio  = iw / ih
    box_ratio  = box_w / box_h

    if img_ratio >= box_ratio:
        new_w = box_w
        new_h = int(box_w / img_ratio)
    else:
        new_h = box_h
        new_w = int(box_h * img_ratio)

    new_left = left + (box_w - new_w) // 2
    new_top  = top  + (box_h - new_h) // 2

    slide.shapes.add_picture(BytesIO(img_bytes), new_left, new_top,
                             width=new_w, height=new_h)


def _fill_name(text_shape, parts):
    tf = text_shape.text_frame

    size = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        size = tf.paragraphs[0].runs[0].font.size or Pt(12)
    else:
        size = Pt(12)

    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    for i, part in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = part
        run.font.size = size
        run.font.bold = i == 0

def generate_pptx(template_filename: str, products: list, output_filename: str):
    prs = Presentation(template_filename)
    prototype_idx = 0
    pair_count = len(_get_name_pic_pairs(prs.slides[prototype_idx]))

    for batch_start in range(0, len(products), pair_count):
        batch = products[batch_start : batch_start + pair_count]
        slide = _duplicate_slide(prs, prototype_idx)

        for prod, (txt_shape, pic_shape) in zip(batch, _get_name_pic_pairs(slide)):
            _fill_name(txt_shape, prod["name"])
            try:
                rsp = requests.get(prod["image"], timeout=10)
                rsp.raise_for_status()
                _replace_picture(slide, pic_shape, rsp.content)
            except Exception as exc:
                print(f"[warn] Couldn’t fetch image {prod['image']}: {exc}")

        extras = _get_name_pic_pairs(slide)[len(batch) :]
        for txt_shape, pic_shape in extras:
            txt_shape.text_frame.text = ""
            slide.shapes._spTree.remove(pic_shape._element)

    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[prototype_idx])

    prs.save(output_filename)
    print(f"Katalog wygenerowany: {output_filename}")
